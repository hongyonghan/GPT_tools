import glob
import os
import random
import re
import string

import gradio as gr

import openai
from icrawler import ImageDownloader
from icrawler.builtin import BingImageCrawler
from uuid import uuid4
from pptx import Presentation

bad_coding_practice = ''.join(random.choice(string.ascii_uppercase + string.ascii_lowercase + string.digits) for _ in
                              range(16))

import gradio as gr

from langchain.document_loaders import PyPDFLoader  # for loading the pdf
from langchain.embeddings import OpenAIEmbeddings  # for creating embeddings
from langchain.vectorstores import Chroma  # for the vectorization part
from langchain.chains import ChatVectorDBChain  # for chatting with the pdf
from langchain.llms import OpenAI  # the LLM model we'll use (CHatGPT)
import shutil
import os

class Chat:
    def __init__(self, pdf, api_input):
        self.api = api_input
        os.environ["OPENAI_API_KEY"] = self.api
        # db = os.getcwd() + "/db"
        # index = os.getcwd() + "/index"
        # print(db)
        # print(index)

        # shutil.rmtree(db)
        # shutil.rmtree(index)

        loader = PyPDFLoader(pdf)
        pages = loader.load_and_split()


        embeddings = OpenAIEmbeddings(openai_api_key=self.api)
        vectordb = Chroma.from_documents(pages, embedding=embeddings, persist_directory="db")
        vectordb.persist()

        self.pdf_qa = ChatVectorDBChain.from_llm(OpenAI(temperature=0.9, model_name="gpt-3.5-turbo",
                                                        openai_api_key=self.api),
                                                 vectordb, return_source_documents=True)

    def question(self, query):
        result = self.pdf_qa({"question": "请使用中文回答" + query, "chat_history": ""})
        print("Answer:")
        print(result["answer"])

        return result["answer"]


def get_result(query):
    openai_api_key = os.getenv("OPENAI_API_KEY")
    embeddings = OpenAIEmbeddings(openai_api_key=openai_api_key)
    vectordb = Chroma(persist_directory='db', embedding_function=embeddings)

    pdf_qa = ChatVectorDBChain.from_llm(OpenAI(temperature=0.9, model_name="gpt-3.5-turbo",
                                                    openai_api_key=openai_api_key),
                                             vectordb, return_source_documents=True)
    result = pdf_qa({"question": "请使用中文回答" + query, "chat_history": ""})

    return result["answer"]


def analyse(pdf_file, api_input):
    session = Chat(pdf_file.name, api_input)
    return session, "文章分析完成"


def ask_question(data, question):
    if data == "":
        return "Please upload PDF file first!"
    return data.question(question)
def refresh_bad_coding_practice():
    global bad_coding_practice
    bad_coding_practice = ''.join(random.choice(string.ascii_uppercase + string.ascii_lowercase + string.digits)
                                  for _ in range(16))
    return
class PrefixNameDownloader(ImageDownloader):

    def get_filename(self, task, default_ext):
        filename = super(PrefixNameDownloader, self).get_filename(
            task, default_ext)
        print(bad_coding_practice)
        return 'prefix_' + bad_coding_practice + filename
def generate_ppt(file, topic, slide_length, api_key):
    print(file.name)

    root = Presentation(file.name)

    openai.api_key = api_key

    message = f"""
    Create content for a slideshow presentation.
    The content's topic is {topic}. 
    The slideshow is {slide_length} slides long. 
    The content is written in the language of the content I give you above.
    
    
    You are allowed to use the following slide types:
    
    Slide types:
    Title Slide - (Title, Subtitle)
    Content Slide - (Title, Content)
    Image Slide - (Title, Content, Image)
    Thanks Slide - (Title)
    
    Put this tag before the Title Slide: [L_TS]
    Put this tag before the Content Slide: [L_CS]
    Put this tag before the Image Slide: [L_IS]
    Put this tag before the Thanks Slide: [L_THS]
    
    Put "[SLIDEBREAK]" after each slide 
    
    For example:
    [L_TS]
    [TITLE]Mental Health[/TITLE]
    
    [SLIDEBREAK]
    
    [L_CS] 
    [TITLE]Mental Health Definition[/TITLE]
    [CONTENT]
    1. Definition: A person’s condition with regard to their psychological and emotional well-being
    2. Can impact one's physical health
    3. Stigmatized too often.
    [/CONTENT]
    
    [SLIDEBREAK]
    
    Put this tag before the Title: [TITLE]
    Put this tag after the Title: [/TITLE]
    Put this tag before the Subitle: [SUBTITLE]
    Put this tag after the Subtitle: [/SUBTITLE]
    Put this tag before the Content: [CONTENT]
    Put this tag after the Content: [/CONTENT]
    Put this tag before the Image: [IMAGE]
    Put this tag after the Image: [/IMAGE]
    
    Elaborate on the Content, provide as much information as possible.
    You put a [/CONTENT] at the end of the Content.
    Do not reply as if you are talking about the slideshow itself. (ex. "Include pictures here about...")
    Do not include any special characters (?, !, ., :, ) in the Title.
    Do not include any additional information in your response and stick to the format."""

    openai.proxy = {'http': "http://127.0.0.1:8001", 'https': 'http://127.0.0.1:8001'}
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "user", "content": message}
        ]
    )

    # """ Ref for slide types:
    # 0 -> title and subtitle
    # 1 -> title and content
    # 2 -> section header
    # 3 -> two content
    # 4 -> Comparison
    # 5 -> Title only
    # 6 -> Blank
    # 7 -> Content with caption
    # 8 -> Pic with caption
    # """

    def delete_all_slides():
        for i in range(len(root.slides) - 1, -1, -1):
            r_id = root.slides._sldIdLst[i].rId
            root.part.drop_rel(r_id)
            del root.slides._sldIdLst[i]

    def create_title_slide(title, subtitle):
        layout = root.slide_layouts[0]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    def create_section_header_slide(title):
        layout = root.slide_layouts[2]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = title

    def create_title_and_content_slide(title, content):
        layout = root.slide_layouts[1]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = content

    def create_title_and_content_and_image_slide(title, content, image_query):
        layout = root.slide_layouts[8]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[2].text = content
        refresh_bad_coding_practice()
        baidu_crawler = BingImageCrawler(downloader_cls=PrefixNameDownloader, storage={'root_dir': os.getcwd()})
        baidu_crawler.crawl(keyword=image_query, max_num=1)
        dir_path = os.path.dirname(os.path.realpath(__file__))
        file_name = glob.glob(f"prefix_{bad_coding_practice}*")
        print(file_name)
        img_path = os.path.join(dir_path, file_name[0])
        slide.shapes.add_picture(img_path, slide.placeholders[1].left, slide.placeholders[1].top,
                                 slide.placeholders[1].width, slide.placeholders[1].height)

    def find_text_in_between_tags(text, start_tag, end_tag):
        start_pos = text.find(start_tag)
        end_pos = text.find(end_tag)
        result = []
        while start_pos > -1 and end_pos > -1:
            text_between_tags = text[start_pos + len(start_tag):end_pos]
            result.append(text_between_tags)
            start_pos = text.find(start_tag, end_pos + len(end_tag))
            end_pos = text.find(end_tag, start_pos)
        res1 = "".join(result)
        res2 = re.sub(r"\[IMAGE\].*?\[/IMAGE\]", '', res1)
        if len(result) > 0:
            return res2
        else:
            return ""

    def search_for_slide_type(text):
        tags = ["[L_TS]", "[L_CS]", "[L_IS]", "[L_THS]"]
        found_text = next((s for s in tags if s in text), None)
        return found_text
    
    def extend_content(slide):
        message =  f"""扩写如下内容，内容200~300字左右：{slide}"""
        
        openai.proxy = {'http': "http://127.0.0.1:8001", 'https': 'http://127.0.0.1:8001'}
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "user", "content": message}
            ]
        )

        print(response['choices'][0]['message']['content'])

        return response['choices'][0]['message']['content']

    def parse_response(reply):
        list_of_slides = reply.split("[SLIDEBREAK]")
        for slide in list_of_slides:
            slide_type = search_for_slide_type(slide)
            # print("1",str(slide))
            if slide_type == "[L_TS]":
                create_title_slide(find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]"),
                                   find_text_in_between_tags(str(slide), "[SUBTITLE]", "[/SUBTITLE]"))
            elif slide_type == "[L_CS]":
                content_ = extend_content(slide)
                
                # create_title_and_content_slide("".join(find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")),
                #                                "".join(find_text_in_between_tags(str(content_), "[CONTENT]",
                #                                                                  "[/CONTENT]")))
                create_title_and_content_slide("".join(find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")),
                                "".join(content_))
            elif slide_type == "[L_IS]":

                create_title_and_content_and_image_slide("".join(find_text_in_between_tags(str(slide), "[TITLE]",
                                                                                           "[/TITLE]")),
                                                         "".join(find_text_in_between_tags(str(slide), "[CONTENT]",
                                                                                           "[/CONTENT]")),
                                                         "".join(find_text_in_between_tags(str(slide), "[IMAGE]",
                                                                                           "[/IMAGE]")))
            elif slide_type == "[L_THS]":
                create_section_header_slide("".join(find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")))

    def find_title():
        return root.slides[0].shapes.title.text

    def get_pdf_content():
        title = get_result("论文的标题")
        author = get_result("论文的作者")
        contribution = get_result("论文的贡献点")
        datasets = get_result("论文的数据集")
        methods = get_result("论文的方法")
        result = get_result("论文的结果")

        title_slide_layout = root.slide_layouts[0]
        slide = root.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = title
        subtitle.text = author

        #create 2nd slide
        bullet_slide_layout1 = root.slide_layouts[1]
        slide2 = root.slides.add_slide(bullet_slide_layout1)
        shapes2 = slide2.shapes
        title_shape2 = shapes2.title
        body_shape2 = shapes2.placeholders[1]
        title_shape2.text = '论文的贡献点'
        tf2 = body_shape2.text_frame
        tf2.text = contribution
        # create 3nd slide
        slide3 = root.slides.add_slide(bullet_slide_layout1)
        shapes3 = slide3.shapes
        title_shape3 = shapes3.title
        body_shape3 = shapes3.placeholders[1]
        title_shape3.text = '论文的数据集'
        tf3 = body_shape3.text_frame
        tf3.text = datasets

        # create 4nd slide
        slide4 = root.slides.add_slide(bullet_slide_layout1)
        shapes4 = slide4.shapes
        title_shape4 = shapes4.title
        body_shape4 = shapes4.placeholders[1]
        title_shape4.text = '论文的方法'
        tf4 = body_shape4.text_frame
        tf4.text = methods

        # create 5nd slide
        slide5 = root.slides.add_slide(bullet_slide_layout1)
        shapes5 = slide5.shapes
        title_shape5 = shapes5.title
        body_shape5 = shapes5.placeholders[1]
        title_shape5.text = '论文的结果'
        tf5 = body_shape4.text_frame
        tf5.text = result

        # create 6nd slide
        title_slide_layout = root.slide_layouts[0]
        slide6 = root.slides.add_slide(title_slide_layout)
        title6 = slide6.shapes.title
        title6.text = "谢谢大家!"



    delete_all_slides()

    parse_response(response['choices'][0]['message']['content'])

    name_ = str(uuid4()).replace('-', '')

    root.save(f"./{name_}.pptx")

    print("done")

    dir_path = "./"
    prefix = "prefix_"

    for file_name in os.listdir(dir_path):
        if file_name.startswith(prefix):
            file_path = os.path.join(dir_path, file_name)
            if os.path.isfile(file_path):
                os.remove(file_path)

    return f"./{name_}.pptx"


# def write_PPT():

with gr.Blocks(title="ChatGPT PPT框架生成") as demo:
    gr.Markdown("""<h1><center>ChatGPT For PPT</center></h1>""")
    data = gr.State()
    with gr.Tab("Upload PDF File"):
        pdf_input = gr.File(label="PDF File")
        api_input = gr.Textbox(label="OpenAI API Key")

        result = gr.Textbox()
        upload_button = gr.Button("Start Analyse")
        question_input = gr.Textbox(label="Your Question", placeholder="Authors of this paper?")
        answer = gr.Textbox(label="Answer")
        ask_button = gr.Button("Ask")

    upload_button.click(fn=analyse, inputs=[pdf_input, api_input], outputs=[data, result])
    ask_button.click(ask_question, inputs=[data, question_input], outputs=answer)
    with gr.Row():
        with gr.Column():
            openai_token = gr.Textbox(label="OpenAI API Key")
            topic = gr.Textbox(label="PPT的主题或内容")
            length = gr.Slider(minimum=1, maximum=20, value=6, label="生成的PPT页数", step=1)
            theme = gr.File(value="./theme.pptx", file_types=['pptx', 'ppt'], label="PPT模版")
            output_file = gr.File(interactive=False)

            topic.submit(
                fn=generate_ppt,
                inputs=[theme, topic, length, openai_token],
                outputs=[output_file]
            )

            submit = gr.Button("生成")
            submit.click(
                fn=generate_ppt,
                inputs=[theme, topic, length, openai_token],
                outputs=[output_file]
            )

if __name__ == "__main__":
    # demo.launch()
    demo.queue().launch(inbrowser=True, show_api=False, server_name="0.0.0.0", server_port=8000)
    demo.close()
